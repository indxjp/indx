"""Extras e2e for the optional PARSER backends, against their real libraries.

Each parser is exercised by a *full pipeline build* (``--parser X``) over a small HTML
file — the one input format all three handle — with the rest of the stack pinned to the
offline core (``hash`` embedder, ``jsonl`` store) so no API key is needed. A real parse
must yield at least one document with non-empty recovered text.

* ``markitdown``   — pure-Python conversion; fast.
* ``docling``      — downloads layout models on first run (``slow``).
* ``unstructured`` — partitions documents locally (``slow``).
"""

from __future__ import annotations

from pathlib import Path

import pytest

from indx import DirectoryPipeline
from indx.parsers.markitdown import MarkItDownParser

pytestmark = pytest.mark.extras


def _build_with_parser(parser: str, corpus: Path) -> None:
    space = DirectoryPipeline(
        seed=0, parser=parser, llm="none", embedder="hash", store="jsonl"
    ).run(corpus)

    assert space.documents(), f"{parser} build produced no documents"
    assert space.chunks, f"{parser} build produced no chunks"
    assert any(c.text.strip() for c in space.chunks), f"{parser} produced only empty text"


def test_markitdown_build(requires_lib, html_corpus: Path) -> None:
    requires_lib("markitdown")
    _build_with_parser("markitdown", html_corpus)


def _make_docx(path: Path) -> bool:
    """Write a tiny real .docx at ``path``; return False if python-docx isn't available."""
    try:
        # python-docx ships as the ``docx`` import name; pulled by markitdown[docx].
        import docx  # type: ignore[import-not-found]
    except ImportError:
        return False
    document = docx.Document()
    document.add_heading("Quarterly Note", level=1)
    document.add_paragraph("Revenue grew across every region.")
    document.add_paragraph("The engineering team shipped the new search ranking model.")
    document.save(str(path))
    return True


def _make_pptx(path: Path) -> bool:
    """Write a tiny real .pptx at ``path``; return False if python-pptx isn't available."""
    try:
        # python-pptx; pulled by markitdown[pptx].
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError:
        return False
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Roadmap"
    slide.placeholders[1].text = "Ship the binary-document parser fix this quarter."
    prs.save(str(path))
    return True


@pytest.mark.parametrize("kind", ["docx", "pptx"])
def test_markitdown_parses_real_binary_office_doc(requires_lib, tmp_path: Path, kind: str) -> None:
    # #20 Bugs 1+2: the whole point of the markitdown parser is converting binary office docs.
    # A real .docx/.pptx (NUL bytes in its zip container, needs the matching markitdown sub-extra)
    # must round-trip through MarkItDownParser into non-empty blocks. This is the test that would
    # have failed on BOTH the over-broad NUL pre-gate (Bug 1) and the missing sub-extra (Bug 2);
    # the HTML-only happy path dodged both. Self-skips when the fixture builder lib is absent.
    requires_lib("markitdown")
    src = tmp_path / f"doc.{kind}"
    made = _make_docx(src) if kind == "docx" else _make_pptx(src)
    if not made:
        pytest.skip(f"fixture builder for .{kind} not installed (python-{kind})")

    parsed = MarkItDownParser().parse(src)

    assert parsed.blocks, f"markitdown recovered no blocks from a real .{kind}"
    assert any(b.text.strip() for b in parsed.blocks), f".{kind} parse produced only empty text"


@pytest.mark.slow
def test_docling_build(requires_lib, html_corpus: Path) -> None:
    requires_lib("docling")
    _build_with_parser("docling", html_corpus)


@pytest.mark.slow
def test_unstructured_build(requires_lib, html_corpus: Path) -> None:
    requires_lib("unstructured")
    _build_with_parser("unstructured", html_corpus)
